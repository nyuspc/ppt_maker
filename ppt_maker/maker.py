"""
Make PowerPoint slides with template and data
@author: Shuo Sun, Pengcheng Song
"""
from pptx import Presentation
from pptx.chart.data import ChartData
import re

class PrsMaker:
    """This class is used to represent a powerpoint slides."""

    def __init__(self, template):
        "init with a template, which is a pptx package slides"
        self.prs = Presentation(template) 

    def __repr__(self):
        return "a ppt_maker maker"

    def text_inject(self, text_map):
        "inject text into template"
        template_tester = r'\[\[(.*?)\]\]'
        for slide in self.prs.slides:
            for shape in slide.shapes:  # 遍历所有shape
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for count in range(len(p.runs)):
                        if re.match(template_tester, p.runs[count].text):
                            key = re.findall(template_tester, p.runs[count].text)
                            key = key[0]
                            p.runs[count].text = text_map[key]
                            print("Replace text {}".format(key))
                            continue
                        if p.runs[count].text == '[[':
                            p.runs[count].text = ''
                            key = p.runs[count+1].text
                            p.runs[count+1].text = text_map[key]
                            p.runs[count+2].text = ''
                            print("Replace text {}".format(key))
    
    def chart_inject(self, chart_data):
        "inject chart data into template"
        for slide in self.prs.slides:
            for shape in slide.shapes: # 遍历所有shape
                if shape.has_chart:
                    chart = shape.chart
                    if chart.has_title:
                        if chart.chart_title.has_text_frame:
                            title = chart.chart_title.text_frame
                            key = title.text
                            try:
                                data = chart_data[key]
                            except:
                                continue
                            for p in title.paragraphs:
                                for run in p.runs:
                                    if run.text == title.text:
                                        run.text = data['title']
                            cd = ChartData()
                            cd.categories = data['categories']
                            for series in data['series']:
                                cd.add_series(series['name'], series['values'])
                            chart.replace_data(cd)
                            print("Replace data of chart {}".format(key))

    def table_inject(self, table_data):
        "inject chart data into template"
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    ncol = len(table.columns)
                    nrow = len(table.rows)
                    for x in range(ncol):
                        for y in range(nrow):
                            tf = table.cell(y,x).text_frame
                            data = table_data[y][x]
                            for p in tf.paragraphs:
                                for run in p.runs:
                                    run.text = data
                                    data = '' #清除其他runs
                    print("Replace data of table")
                    
    def save(self, path):
        "save pptx file"
        self.prs.save(path)

